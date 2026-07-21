import os
import sys

def build_presentation():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        print("python-pptx is not installed yet.")
        return False

    prs = Presentation()
    # Use 16:9 widescreen slides
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color Palette Constants
    COLOR_BG = RGBColor(13, 13, 13)       # #0D0D0D dark surface
    COLOR_BORDER = RGBColor(30, 30, 30)   # #1E1E1E border dark grey
    COLOR_GREEN = RGBColor(48, 209, 88)    # #30D158 premium neon green
    COLOR_WHITE = RGBColor(255, 255, 255)  # #FFFFFF primary text
    COLOR_MUTED = RGBColor(140, 140, 140)  # #8C8C8C secondary/muted text

    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

    def add_header(slide, title_text, category_text="LOG!T ARCHITECTURE"):
        # Header category tracker
        tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        tf_cat = tx_cat.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Inter"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_GREEN

        # Main slide title
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Poppins"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Dark, Minimalist)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1)

    # Accent decorative box
    shape = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.12), Inches(3.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_GREEN
    shape.line.fill.background()

    # Title & Subtitle text box
    tx_box = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(3.5))
    tf = tx_box.text_frame
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = "LOG!T"
    p1.font.name = "Poppins"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "Technical Architecture & System Design Specification"
    p2.font.name = "Inter"
    p2.font.size = Pt(22)
    p2.font.color.rgb = COLOR_GREEN
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "A premium, cloud-first movie tracking app built on vanilla HTML, CSS, and JS"
    p3.font.name = "Inter"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_MUTED

    # ----------------------------------------------------
    # SLIDE 2: Product Overview
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2)
    add_header(slide2, "System Concept & User Capabilities", "Product Overview")

    # Column 1: Core Goal
    tx_col1 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    tf1 = tx_col1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Core Value Proposition"
    p.font.name = "Poppins"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(14)

    bullet_points = [
        "Search, rate, and track films in a streamlined personal library.",
        "Zero setup overhead - completely serverless static application.",
        "Beautiful, unified dark mode theme optimized for movie collection scanning.",
        "Comprehensive dashboard listing statistics, runtime charts, and watch schedules."
    ]
    for pt in bullet_points:
        p_pt = tf1.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Inter"
        p_pt.font.size = Pt(14)
        p_pt.font.color.rgb = COLOR_MUTED
        p_pt.space_after = Pt(10)

    # Column 2: Architecture Highlights
    tx_col2 = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.8))
    tf2 = tx_col2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Architectural Pillars"
    p.font.name = "Poppins"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(14)

    highlights = [
        "Cloud Sync — Synced in Supabase database instantly with low-latency indexing.",
        "Auto Backup — Seamlessly triggers RFC 2046 multipart backup payloads to Google Drive.",
        "Vanilla Core — Zero bundlers, compile phases, or external package bloating.",
        "Missing Data Auditing — Visual indicator indicators (red dots) for items missing TMDB ratings, runtime, or overview info."
    ]
    for hl in highlights:
        p_hl = tf2.add_paragraph()
        p_hl.text = "✓ " + hl
        p_hl.font.name = "Inter"
        p_hl.font.size = Pt(14)
        p_hl.font.color.rgb = COLOR_MUTED
        p_hl.space_after = Pt(10)

    # ----------------------------------------------------
    # SLIDE 3: Tech Stack
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3)
    add_header(slide3, "Technology Stack & API Integrations", "Tech Stack")

    # We will draw 3 boxes for the core stack items
    boxes = [
        {
            "title": "Front End & Style",
            "subtitle": "Vanilla Core",
            "body": "• Semantic HTML5 structure\n• Custom CSS variables\n• Responsive layout grids\n• Vanilla ES6 modules\n• Cache-controlled assets",
            "left": Inches(0.8), "width": Inches(3.6)
        },
        {
            "title": "Database & Storage",
            "subtitle": "Supabase Cloud DB",
            "body": "• PostgreSQL database\n• Supabase Auth engine\n• Row-Level Security (RLS)\n• Real-time updates\n• Local Storage caching",
            "left": Inches(4.8), "width": Inches(3.6)
        },
        {
            "title": "External APIs",
            "subtitle": "TMDB & Google",
            "body": "• TMDB Search & Detail API\n• Google Identity Services (GIS)\n• Drive API v3 integrations\n• OAuth2 Client tokens\n• RFC 2046 multipart uploads",
            "left": Inches(8.8), "width": Inches(3.6)
        }
    ]

    for b in boxes:
        # Background box
        box_shape = slide3.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, b["left"], Inches(2.0), b["width"], Inches(4.5)
        )
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = COLOR_BG
        box_shape.line.color.rgb = COLOR_BORDER
        box_shape.line.width = Pt(1.5)

        # Text Frame inside box
        tb = slide3.shapes.add_textbox(b["left"] + Inches(0.2), Inches(2.2), b["width"] - Inches(0.4), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = b["title"]
        p_t.font.name = "Poppins"
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_WHITE
        p_t.space_after = Pt(4)

        p_s = tf.add_paragraph()
        p_s.text = b["subtitle"]
        p_s.font.name = "Inter"
        p_s.font.size = Pt(12)
        p_s.font.bold = True
        p_s.font.color.rgb = COLOR_GREEN
        p_s.space_after = Pt(14)

        p_b = tf.add_paragraph()
        p_b.text = b["body"]
        p_b.font.name = "Inter"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = COLOR_MUTED
        p_b.space_after = Pt(6)

    # ----------------------------------------------------
    # SLIDE 4: Google Drive Integration Overhaul (v3.3.0)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4)
    add_header(slide4, "Re-engineered Google Drive API Core", "Google Drive Sync")

    tx_drv1 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tfd1 = tx_drv1.text_frame
    tfd1.word_wrap = True

    p = tfd1.paragraphs[0]
    p.text = "Modern Fetch Architecture"
    p.font.name = "Poppins"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(12)

    drive_points = [
        "Replaced legacy XMLHttpRequests with clean async/await fetch() functions.",
        "Built centralized _apiFetch() wrapper supporting authorization headers injection and token auto-refresh.",
        "Smart OAuth Token Lifetime: Automatically detects expired tokens (401 errors), wipes local caches, and updates the connection state UI.",
        "Comprehensive user profiling - calls UserInfo endpoint or falls back to Google Drive About API."
    ]
    for pt in drive_points:
        p_pt = tfd1.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Inter"
        p_pt.font.size = Pt(12.5)
        p_pt.font.color.rgb = COLOR_MUTED
        p_pt.space_after = Pt(8)

    tx_drv2 = slide4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tfd2 = tx_drv2.text_frame
    tfd2.word_wrap = True

    p = tfd2.paragraphs[0]
    p.text = "Payload Optimization & Folders"
    p.font.name = "Poppins"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(12)

    drive_payloads = [
        "RFC 2046 Multipart Request Body: Uploads files using multipart/related boundaries to pack binary metadata and file data in one request.",
        "Folder isolation: Scans and isolates a dedicated root-level folder named 'Logit' on Google Drive.",
        "Automatic file naming: Generates files dynamically (e.g., logit-{count}-movies-{timestamp}.json).",
        "Targeted Queries: Restricts folder search using highly specific Google Query DSL filters to exclude trashed files."
    ]
    for pl in drive_payloads:
        p_pl = tfd2.add_paragraph()
        p_pl.text = "✓ " + pl
        p_pl.font.name = "Inter"
        p_pl.font.size = Pt(12.5)
        p_pl.font.color.rgb = COLOR_MUTED
        p_pl.space_after = Pt(8)

    # ----------------------------------------------------
    # SLIDE 5: User Interface Layout & Spacing Design
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5)
    add_header(slide5, "Refined Minimal Layout & Spacing", "UI/UX & Spacing")

    # Column 1
    tx_ui1 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tfu1 = tx_ui1.text_frame
    tfu1.word_wrap = True

    p = tfu1.paragraphs[0]
    p.text = "Apple-Style Minimalism"
    p.font.name = "Poppins"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(12)

    ui_points = [
        "Unified Cloud Card — Displays Google Drive connection state, status badge, and email in a compact, unified pill card.",
        "Clean Segmented Settings — Grouped settings into 4 semantic blocks (Account, Google Drive, Backup & Data, Danger Zone).",
        "Perfect Title Centering — Positioned the profile edit pen absolutely so it does not affect name text alignment flow.",
        "Interactive Status Badge — Pulsing neon status dot dynamically highlights connection state changes."
    ]
    for pt in ui_points:
        p_pt = tfu1.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Inter"
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = COLOR_MUTED
        p_pt.space_after = Pt(8)

    # Column 2
    tx_ui2 = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tfu2 = tx_ui2.text_frame
    tfu2.word_wrap = True

    p = tfu2.paragraphs[0]
    p.text = "Layout Grid & Media Queries"
    p.font.name = "Poppins"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_after = Pt(12)

    grid_points = [
        "Consistent Vertical Spacing — Verified content block spacing remains exactly 24px across both desktop and mobile screens.",
        "Adaptive Breakpoints — Overhauled media queries to wrap lists, buttons, and badges dynamically without overflow.",
        "Chevron-Aligned Items — Flexbox list items feature right-aligned indicator chevrons for clean scanning.",
        "Mobile Contrast Optimization — Responsive dark theme uses calibrated grey borders (#1e1e1e) to separate interactive buttons."
    ]
    for gp in grid_points:
        p_gp = tfu2.add_paragraph()
        p_gp.text = "✓ " + gp
        p_gp.font.name = "Inter"
        p_gp.font.size = Pt(13)
        p_gp.font.color.rgb = COLOR_MUTED
        p_gp.space_after = Pt(8)

    # Save presentation
    filename = "Logit_System_Architecture_Presentation.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")
    return True

if __name__ == "__main__":
    build_presentation()
